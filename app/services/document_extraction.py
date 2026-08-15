from dataclasses import dataclass
from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader


@dataclass
class TextChunk:
    text: str
    chunk_index: int
    source_document: str
    page_number: int | None = None


def extract_text_from_pdf(path: str) -> str:
    """Extract text page by page using pdfplumber, falling back to pypdf on failure."""
    try:
        with pdfplumber.open(path) as pdf:
            pages = [page.extract_text() or "" for page in pdf.pages]
        return "\n\n".join(pages)
    except Exception:
        reader = PdfReader(path)
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages)


def extract_text_from_docx(path: str) -> str:
    doc = DocxDocument(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_from_pptx(path: str) -> str:
    """Extract slide title/body text and speaker notes, in slide order."""
    prs = Presentation(path)
    slides_text: list[str] = []

    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            elif shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Notes: {notes}")
        slides_text.append("\n".join(parts))

    return "\n\n".join(s for s in slides_text if s)


def extract_text(path: str) -> str:
    suffix = Path(path).suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(path)
    if suffix == ".docx":
        return extract_text_from_docx(path)
    if suffix == ".pptx":
        return extract_text_from_pptx(path)
    if suffix == ".txt":
        return Path(path).read_text(encoding="utf-8")
    if suffix == ".ppt":
        raise ValueError(
            "Legacy .ppt is not directly supported — convert to .pptx first "
            "(e.g. open and Save As in PowerPoint)."
        )
    raise ValueError(f"Unsupported file type: {suffix}")


def chunk_text(
    text: str,
    source_document: str,
    chunk_size: int = 800,
    overlap: int = 150,
) -> list[TextChunk]:
    """Split text into overlapping word-based chunks for embedding/retrieval."""
    words = text.split()
    if not words:
        return []

    chunks: list[TextChunk] = []
    start = 0
    index = 0
    step = max(chunk_size - overlap, 1)

    while start < len(words):
        end = start + chunk_size
        chunk_words = words[start:end]
        chunks.append(
            TextChunk(
                text=" ".join(chunk_words),
                chunk_index=index,
                source_document=source_document,
            )
        )
        index += 1
        start += step

    return chunks


def extract_and_chunk(path: str, chunk_size: int = 800, overlap: int = 150) -> list[TextChunk]:
    """Small overlapping chunks, sized for embedding/RAG retrieval (Phase 3)."""
    text = extract_text(path)
    return chunk_text(text, source_document=Path(path).name, chunk_size=chunk_size, overlap=overlap)


def extract_and_chunk_for_extraction(
    path: str, chunk_size: int = 4000, overlap: int = 0
) -> list[TextChunk]:
    """Large, mostly non-overlapping chunks for LLM graph extraction (Phase 2.2).

    Bigger chunks mean fewer LLM calls per document and let the model see
    relationships that span what would otherwise be separate retrieval chunks.
    """
    text = extract_text(path)
    return chunk_text(text, source_document=Path(path).name, chunk_size=chunk_size, overlap=overlap)
